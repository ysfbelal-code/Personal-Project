import http.server
import json
import urllib.request
import urllib.error
import urllib.parse
import os

PORT = 8080
GROQ_URL = "https://api.groq.com/openai/v1/chat/completions"
HTML_FILE = os.path.join(os.path.dirname(__file__), "index.html")


class ElixirHandler(http.server.BaseHTTPRequestHandler):

    def log_message(self, fmt, *args):
        print(f"  {self.address_string()} → {fmt % args}")

    # ── GET: serve elixir.html ───────────────────────────────────────────────
    def do_GET(self):
        try:
            with open(HTML_FILE, "rb") as f:
                data = f.read()
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.send_header("Content-Length", str(len(data)))
            self.end_headers()
            self.wfile.write(data)
        except FileNotFoundError:
            self.send_error(404, "elixir.html not found next to server.py")

    # ── OPTIONS: CORS preflight ──────────────────────────────────────────────
    def do_OPTIONS(self):
        self.send_response(204)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, GET, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.end_headers()

    # ── POST: route requests ─────────────────────────────────────────────────
    def do_POST(self):
        parsed = urllib.parse.urlparse(self.path)
        route  = parsed.path

        if route not in ("/groq", "/models", "/extract"):
            self.send_error(404)
            return

        if route == "/extract":
            self._handle_extract(parsed)
            return

        # Groq routes — require JSON body with API key
        length = int(self.headers.get("Content-Length", 0))
        raw    = self.rfile.read(length)
        try:
            payload = json.loads(raw)
        except json.JSONDecodeError:
            self.send_error(400, "Invalid JSON")
            return

        api_key = payload.pop("key", "")
        if not api_key:
            self._json_error(400, "No API key provided")
            return

        if route == "/models":
            print(f"  [models] Using API key: {api_key[:20]}...")
            print(f"  [models] Requesting: https://api.groq.com/openai/v1/models")
            req = urllib.request.Request(
                "https://api.groq.com/openai/v1/models",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "User-Agent": "Mozilla/5.0",
                },
                method="GET",
            )
        else:
            req = urllib.request.Request(
                GROQ_URL,
                data=json.dumps(payload).encode(),
                headers={
                    "Authorization":   f"Bearer {api_key}",
                    "Content-Type":    "application/json",
                    "User-Agent":      "Mozilla/5.0",
                    "Accept":          "application/json",
                    "Accept-Language": "en-US,en;q=0.9",
                },
                method="POST",
            )

        MAX_RETRIES = 3
        last_err = None
        for attempt in range(MAX_RETRIES):
            try:
                with urllib.request.urlopen(req, timeout=60) as resp:
                    body = resp.read()
                    print(f"  [groq] Success: {resp.status}")
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.send_header("Access-Control-Allow-Origin", "*")
                self.end_headers()
                self.wfile.write(body)
                return
            except urllib.error.HTTPError as e:
                body = e.read()
                print(f"  [groq] HTTP Error {e.code}: {body[:200]}")
                
                # 405 error from /models endpoint? Return hardcoded Groq models as fallback
                if e.code == 405 and route == "/models":
                    print(f"  [models] /models endpoint returned 405, using fallback model list")
                    fallback = {
                        "object": "list",
                        "data": [
                            {"id": "mixtral-8x7b-32768", "object": "model", "owned_by": "groq"},
                            {"id": "llama-3.1-70b-versatile", "object": "model", "owned_by": "groq"},
                            {"id": "llama-3.1-8b-instant", "object": "model", "owned_by": "groq"},
                            {"id": "llama-3.2-11b-vision-preview", "object": "model", "owned_by": "groq"},
                            {"id": "llama-3.2-1b-preview", "object": "model", "owned_by": "groq"},
                        ]
                    }
                    body = json.dumps(fallback).encode()
                
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.send_header("Access-Control-Allow-Origin", "*")
                self.end_headers()
                self.wfile.write(body)
                return
            except urllib.error.URLError as e:
                last_err = e
                print(f"  [groq] URL Error (attempt {attempt+1}): {e.reason}")
                if attempt < MAX_RETRIES - 1:
                    import time; time.sleep(1.5 * (attempt + 1))

        self._json_error(502,
            f"Could not reach Groq after {MAX_RETRIES} attempts: {last_err.reason}")

    # ────────────────────────────────────────────────────────────────────────
    # /extract  — receive raw file bytes, return extracted text as JSON
    # ────────────────────────────────────────────────────────────────────────
    def _handle_extract(self, parsed):
        qs       = urllib.parse.parse_qs(parsed.query)
        filename = urllib.parse.unquote(qs.get("filename", ["unknown"])[0])

        # Read body — honour Content-Length or fall back to chunked read
        length = int(self.headers.get("Content-Length", 0))
        if length > 0:
            raw = self.rfile.read(length)
        else:
            chunks = []
            while True:
                chunk = self.rfile.read(65536)
                if not chunk:
                    break
                chunks.append(chunk)
            raw = b"".join(chunks)

        print(f"  [extract] {filename} ({len(raw):,} bytes)")

        if not raw:
            self._json_resp({"text": "", "error":
                "Empty file received — no bytes arrived. Try uploading again."})
            return

        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        try:
            if   ext == "txt":               text = self._extract_txt(raw)
            elif ext == "pptx":              text = self._extract_pptx(raw)
            elif ext == "ppt":               text = self._extract_ppt(raw)
            elif ext == "docx":              text = self._extract_docx(raw)
            elif ext == "doc":               text = self._extract_doc(raw)
            elif ext == "pdf":               text = self._extract_pdf(raw)
            elif ext == "xlsx":              text = self._extract_xlsx(raw)
            elif ext == "xls":               text = self._extract_xls(raw)
            elif ext == "rtf":               text = self._extract_rtf(raw)
            elif ext in ("odt","odp","ods"): text = self._extract_odf(raw)
            else:
                self._json_resp({"text": "", "error":
                    f"Unsupported file type '.{ext or 'unknown'}'. "
                    "Supported: txt, pptx, ppt, docx, doc, pdf, xlsx, xls, rtf, odt, odp, ods."})
                return

            print(f"  [extract] extracted {len(text):,} chars from {filename}")
            self._json_resp({"text": text, "error": None})

        except Exception as e:
            msg = f"{type(e).__name__}: {e}"
            print(f"  [extract] ERROR: {msg}")
            self._json_resp({"text": "", "error": msg})

    # ════════════════════════════════════════════════════════════════════════
    # Extractors
    # ════════════════════════════════════════════════════════════════════════

    # ── .txt ─────────────────────────────────────────────────────────────────
    def _extract_txt(self, raw):
        """Plain text file — decode with UTF-8, fall back to latin-1."""
        try:
            return raw.decode("utf-8", errors="replace").strip()
        except Exception:
            return raw.decode("latin-1", errors="replace").strip()

    # ── .pptx ────────────────────────────────────────────────────────────────
    def _extract_pptx(self, raw):
        """Modern PowerPoint — bold/underline tagging with para.text fallback."""
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError(
                "python-pptx not installed. Run: python3 -m pip install python-pptx")
        import io
        prs   = Presentation(io.BytesIO(raw))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    # Try runs first to capture bold/underline markers
                    parts = []
                    for run in para.runs:
                        t = run.text
                        if not t:
                            continue
                        b = bool(run.font.bold)
                        u = bool(run.font.underline)
                        if b and u:   parts.append(f"[KEY TERM: {t}]")
                        elif b:       parts.append(f"[BOLD: {t}]")
                        elif u:       parts.append(f"[UNDERLINE: {t}]")
                        else:         parts.append(t)
                    line = "".join(parts).strip()
                    # Fallback: para.text captures text with no explicit runs
                    if not line:
                        line = para.text.strip()
                    if line:
                        lines.append(line)
        return "\n".join(lines)

    # ── .ppt (OLE2 binary) ───────────────────────────────────────────────────
    def _extract_ppt(self, raw):
        """
        Legacy PowerPoint 97-2003 (.ppt) OLE2 binary format.

        PPT records have an 8-byte header:
          [0:2]  verAndInstance — if low nibble == 0xF, it's a CONTAINER
          [2:4]  record type
          [4:8]  record length (bytes of content, NOT including the header)

        BUG IN PREVIOUS VERSION: always did offset = data_e, which skipped
        over every container record AND all its nested children. Text atoms
        live INSIDE container records, so they were never reached.

        FIX: for container records, step PAST only the 8-byte header (enter
        the container). For atom records, step past header + content as before.
        The nested records' cumulative lengths naturally bring the offset to
        exactly the end of the parent container — no stack needed.
        """
        try:
            import olefile
        except ImportError:
            raise RuntimeError(
                "olefile not installed. Run: python3 -m pip install olefile")
        import io, struct

        ole = olefile.OleFileIO(io.BytesIO(raw))
        if not ole.exists("PowerPoint Document"):
            ole.close()
            raise ValueError(
                "Not a valid .ppt file — 'PowerPoint Document' stream missing.")
        stream = ole.openstream("PowerPoint Document").read()
        ole.close()

        # Record type constants
        RT_Slide         = 0x03EE   # container: marks a new slide
        RT_TextCharsAtom = 0x0FA0   # atom: UTF-16-LE text
        RT_TextBytesAtom = 0x0FA8   # atom: Latin-1 text

        lines     = []
        slide_num = 0
        offset    = 0

        while offset + 8 <= len(stream):
            ver_inst = struct.unpack_from("<H", stream, offset)[0]
            rec_type = struct.unpack_from("<H", stream, offset + 2)[0]
            rec_len  = struct.unpack_from("<I", stream, offset + 4)[0]
            is_container = (ver_inst & 0x0F) == 0x0F

            if is_container:
                # Label slides; then ENTER the container (skip only the header)
                if rec_type == RT_Slide:
                    slide_num += 1
                    lines.append(f"\n=== Slide {slide_num} ===")
                offset += 8  # step into container — process nested records next
            else:
                # Atom record — process content then skip past it
                data_s = offset + 8
                data_e = data_s + rec_len
                if data_e > len(stream):
                    break

                if rec_type == RT_TextCharsAtom and rec_len >= 2:
                    try:
                        t = stream[data_s:data_e].decode(
                            "utf-16-le", errors="replace").strip()
                        if t:
                            lines.append(t)
                    except Exception:
                        pass

                elif rec_type == RT_TextBytesAtom and rec_len >= 1:
                    try:
                        t = stream[data_s:data_e].decode(
                            "latin-1", errors="replace").strip()
                        if t:
                            lines.append(t)
                    except Exception:
                        pass

                offset = data_e

        return "\n".join(lines)

    # ── .docx ────────────────────────────────────────────────────────────────
    def _extract_docx(self, raw):
        """Modern Word — bold/underline tagging with para.text fallback."""
        try:
            from docx import Document
        except ImportError:
            raise RuntimeError(
                "python-docx not installed. Run: python3 -m pip install python-docx")
        import io
        doc   = Document(io.BytesIO(raw))
        lines = []
        for para in doc.paragraphs:
            parts = []
            for run in para.runs:
                t = run.text
                if not t:
                    continue
                b = bool(run.bold)
                u = bool(run.underline)
                if b and u:   parts.append(f"[KEY TERM: {t}]")
                elif b:       parts.append(f"[BOLD: {t}]")
                elif u:       parts.append(f"[UNDERLINE: {t}]")
                else:         parts.append(t)
            line = "".join(parts).strip()
            # Fallback for paragraphs with no explicit runs
            if not line:
                line = para.text.strip()
            if line:
                lines.append(line)
        return "\n".join(lines)

    # ── .doc (OLE2 binary) ───────────────────────────────────────────────────
    def _extract_doc(self, raw):
        """
        Legacy Word 97-2003 (.doc) binary format.

        BUG IN PREVIOUS VERSION: scanned the WordDocument stream from byte 0
        as raw UTF-16-LE, which reads the binary FIB header as text — producing
        garbage that filters to nothing.

        FIX: scan the ENTIRE raw file bytes for runs of printable ASCII
        characters encoded as UTF-16-LE (each char is [0x20-0x7E] + 0x00).
        This reliably finds all readable text regardless of offset or structure.
        Filter by minimum run length and minimum printable ratio.
        """
        import re

        # Match 8+ consecutive printable UTF-16-LE ASCII characters
        pattern = re.compile(b"(?:[\x20-\x7e][\x00]){8,}")
        seen  = set()
        lines = []

        for match in pattern.finditer(raw):
            try:
                text = match.group().decode("utf-16-le", errors="replace").strip()
            except Exception:
                continue
            if not text or text in seen:
                continue
            # Must be ≥70% printable to filter out binary noise
            printable = sum(1 for c in text if 0x20 <= ord(c) <= 0x7E)
            if printable / len(text) >= 0.70:
                seen.add(text)
                lines.append(text)

        return "\n".join(lines)

    # ── .pdf ─────────────────────────────────────────────────────────────────
    def _extract_pdf(self, raw):
        """PDF — tries pypdf then PyPDF2."""
        import io
        for mod_name in ("pypdf", "PyPDF2"):
            try:
                mod    = __import__(mod_name)
                reader = mod.PdfReader(io.BytesIO(raw))
                pages  = []
                for i, page in enumerate(reader.pages, 1):
                    txt = (page.extract_text() or "").strip()
                    if txt:
                        pages.append(f"=== Page {i} ===\n{txt}")
                return "\n".join(pages)
            except ImportError:
                continue
        raise RuntimeError("pypdf not installed. Run: python3 -m pip install pypdf")

    # ── .xlsx ────────────────────────────────────────────────────────────────
    def _extract_xlsx(self, raw):
        """Modern Excel."""
        try:
            import openpyxl
        except ImportError:
            raise RuntimeError(
                "openpyxl not installed. Run: python3 -m pip install openpyxl")
        import io
        wb    = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows():
                vals = [str(c.value) for c in row if c.value is not None]
                if vals:
                    lines.append("\t".join(vals))
        return "\n".join(lines)

    # ── .xls (OLE2 binary) ───────────────────────────────────────────────────
    def _extract_xls(self, raw):
        """Legacy Excel 97-2003 via xlrd."""
        try:
            import xlrd
        except ImportError:
            raise RuntimeError(
                "xlrd not installed. Run: python3 -m pip install xlrd")
        wb    = xlrd.open_workbook(file_contents=raw)
        lines = []
        for sheet in wb.sheets():
            lines.append(f"=== Sheet: {sheet.name} ===")
            for row_idx in range(sheet.nrows):
                vals = [str(sheet.cell_value(row_idx, c)).strip()
                        for c in range(sheet.ncols)
                        if str(sheet.cell_value(row_idx, c)).strip()]
                if vals:
                    lines.append("\t".join(vals))
        return "\n".join(lines)

    # ── .rtf ─────────────────────────────────────────────────────────────────
    def _extract_rtf(self, raw):
        """Rich Text Format — pure Python regex stripper, no library needed."""
        import re
        text = raw.decode("latin-1", errors="replace")
        text = re.sub(r"\{\\[^{}]{0,32}\}", "", text)   # remove isolated control groups
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text) # remove control words
        text = re.sub(r"\\[{}\\]", "", text)             # remove escaped braces/backslash
        text = text.replace("{", "").replace("}", "")
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        return "\n".join(lines)

    # ── .odt / .odp / .ods ───────────────────────────────────────────────────
    def _extract_odf(self, raw):
        """OpenDocument formats — unzip and strip content.xml."""
        import io, zipfile, re
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            if "content.xml" not in zf.namelist():
                raise ValueError("Not a valid ODF file — content.xml missing.")
            xml = zf.read("content.xml").decode("utf-8", errors="replace")
        xml = re.sub(r"<(?:text:p|text:h|table:table-row)[^>]*>", "\n", xml)
        xml = re.sub(r"<table:table-cell[^>]*>", "\t", xml)
        xml = re.sub(r"<[^>]+>", "", xml)
        for ent, ch in (("&amp;","&"),("&lt;","<"),("&gt;",">"),
                        ("&quot;",'"'),("&apos;","'")):
            xml = xml.replace(ent, ch)
        lines = [l.strip() for l in xml.splitlines() if l.strip()]
        return "\n".join(lines)

    # ── helpers ──────────────────────────────────────────────────────────────
    def _json_resp(self, data):
        body = json.dumps(data).encode()
        self.send_response(200)
        self.send_header("Content-Type", "application/json")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.end_headers()
        self.wfile.write(body)

    def _json_error(self, code, msg):
        body = json.dumps({"error": {"message": msg}}).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.end_headers()
        self.wfile.write(body)


# ── Boot ─────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import socket
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(("8.8.8.8", 80))
        lan_ip = s.getsockname()[0]
        s.close()
    except Exception:
        lan_ip = "unknown"

    print(f"\n  Elixir is running!\n")
    print(f"  Computer : http://localhost:{PORT}")
    print(f"  Phone    : http://{lan_ip}:{PORT}  (must be on same WiFi)\n")
    print("  Press Ctrl+C to stop.\n")

    server = http.server.HTTPServer(("", PORT), ElixirHandler)
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\n  Server stopped.")
